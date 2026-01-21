"""
GÉNÉRATION PRÉSENTATION POWERPOINT
Crée une présentation complète pour l'opérationnalisation du modèle

Usage:
    python ml_pipeline_v2/generate_powerpoint.py --output presentation_scoring.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path
import argparse


class PresentationBuilder:
    """Constructeur de présentation PowerPoint"""

    def __init__(self, output_file):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        self.output_file = output_file

        # Couleurs
        self.COLOR_TITLE = RGBColor(44, 62, 80)  # Bleu foncé
        self.COLOR_SUBTITLE = RGBColor(52, 73, 94)
        self.COLOR_ACCENT = RGBColor(46, 204, 113)  # Vert
        self.COLOR_WARNING = RGBColor(231, 76, 60)  # Rouge

    def add_title_slide(self):
        """Slide 1: Page de titre"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Titre principal
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "OPÉRATIONNALISATION DU MODÈLE"
        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_TITLE

        # Sous-titre
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(8), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Scoring Automatisé des Réclamations Financières"
        p = subtitle_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(28)
        p.font.color.rgb = self.COLOR_SUBTITLE

        # Date
        date_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
        date_frame = date_box.text_frame
        from datetime import datetime
        date_frame.text = f"{datetime.now().strftime('%B %Y')}"
        p = date_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.color.rgb = self.COLOR_SUBTITLE

        print("✅ Slide 1: Page de titre")

    def add_agenda_slide(self):
        """Slide 2: Agenda"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])  # Title and Content

        # Titre
        title = slide.shapes.title
        title.text = "AGENDA"
        title.text_frame.paragraphs[0].font.size = Pt(40)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.COLOR_TITLE

        # Contenu
        content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        agenda_items = [
            ("I.", "ÉTAT DES LIEUX", [
                "Évolution du volume et montant des réclamations",
                "Analyse fondée vs non fondée",
                "Répartitions par famille et marché"
            ]),
            ("II.", "PRÉSENTATION DU MODÈLE", [
                "Architecture (3 piliers)",
                "Couche analytique (IA)",
                "Couche décisionnelle (règles métier)"
            ]),
            ("III.", "RÉSULTATS 2025 & GAINS", [
                "Performance du modèle",
                "Calcul du gain financier et temps",
                "ROI et recommandations"
            ])
        ]

        for num, section, items in agenda_items:
            # Section principale
            p = tf.add_paragraph()
            p.text = f"{num} {section}"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = self.COLOR_ACCENT
            p.level = 0
            p.space_after = Pt(6)

            # Sous-items
            for item in items:
                p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(16)
                p.font.color.rgb = self.COLOR_SUBTITLE
                p.level = 1
                p.space_after = Pt(3)

            # Espace après section
            p = tf.add_paragraph()
            p.text = ""
            p.space_after = Pt(10)

        print("✅ Slide 2: Agenda")

    def add_section_slide(self, section_number, section_title):
        """Slide de séparation de section"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Fond coloré
        background = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(10), Inches(7.5)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLOR_ACCENT
        background.line.fill.background()

        # Numéro de section
        num_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        num_frame = num_box.text_frame
        num_frame.text = section_number
        p = num_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(80)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Titre de section
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = section_title
        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        print(f"✅ Slide: Section {section_number} - {section_title}")

    def add_evolution_volume_slide(self):
        """Slide 3: Évolution volume et montant"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Titre
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        tf = title_box.text_frame
        tf.text = "ÉVOLUTION DES RÉCLAMATIONS (2023-2025)"
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_TITLE
        p.alignment = PP_ALIGN.CENTER

        # Placeholder pour graphique
        placeholder = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.8), Inches(1.2),
            Inches(8.4), Inches(5.5)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
        placeholder.line.color.rgb = RGBColor(200, 200, 200)

        # Texte dans placeholder
        text_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8.4), Inches(1))
        tf = text_box.text_frame
        tf.text = "📊 INSÉRER ICI: 01_evolution_volume_montant.png"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(100, 100, 100)

        # Note
        note_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.4))
        tf = note_box.text_frame
        tf.text = "Points clés: • Tendance du volume  • Évolution du montant total  • Taux de croissance"
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = RGBColor(100, 100, 100)

        print("✅ Slide 3: Évolution volume et montant")

    def add_fondee_vs_non_fondee_slide(self):
        """Slide 4: Fondée vs Non fondée"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Titre
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        tf = title_box.text_frame
        tf.text = "ANALYSE FONDÉE vs NON FONDÉE"
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_TITLE
        p.alignment = PP_ALIGN.CENTER

        # Placeholder pour graphique
        placeholder = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.8), Inches(1.2),
            Inches(8.4), Inches(5.5)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
        placeholder.line.color.rgb = RGBColor(200, 200, 200)

        # Texte dans placeholder
        text_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8.4), Inches(1))
        tf = text_box.text_frame
        tf.text = "📊 INSÉRER ICI: 02_fondee_vs_non_fondee.png"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(100, 100, 100)

        # Note
        note_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.4))
        tf = note_box.text_frame
        tf.text = "Points clés: • % fondée par année  • Montants par catégorie  • Évolution du taux de fondée"
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = RGBColor(100, 100, 100)

        print("✅ Slide 4: Fondée vs Non fondée")

    def add_repartition_famille_slide(self):
        """Slide 5: Répartition par famille"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Titre
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        tf = title_box.text_frame
        tf.text = "RÉPARTITION PAR FAMILLE DE PRODUIT"
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_TITLE
        p.alignment = PP_ALIGN.CENTER

        # Placeholder pour graphique
        placeholder = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.8), Inches(1.2),
            Inches(8.4), Inches(5.5)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
        placeholder.line.color.rgb = RGBColor(200, 200, 200)

        # Texte dans placeholder
        text_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8.4), Inches(1))
        tf = text_box.text_frame
        tf.text = "📊 INSÉRER ICI: 03_repartition_famille.png"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(100, 100, 100)

        # Note
        note_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.4))
        tf = note_box.text_frame
        tf.text = "Points clés: • Top 5 familles par année  • Distribution en nombre vs montant"
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = RGBColor(100, 100, 100)

        print("✅ Slide 5: Répartition par famille")

    def add_repartition_marche_slide(self):
        """Slide 6: Répartition par marché"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Titre
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        tf = title_box.text_frame
        tf.text = "RÉPARTITION PAR MARCHÉ"
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_TITLE
        p.alignment = PP_ALIGN.CENTER

        # Placeholder pour graphique
        placeholder = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.8), Inches(1.2),
            Inches(8.4), Inches(5.5)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
        placeholder.line.color.rgb = RGBColor(200, 200, 200)

        # Texte dans placeholder
        text_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8.4), Inches(1))
        tf = text_box.text_frame
        tf.text = "📊 INSÉRER ICI: 04_repartition_marche.png"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(100, 100, 100)

        # Note
        note_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.4))
        tf = note_box.text_frame
        tf.text = "Points clés: • Distribution par segment de marché  • Comparaison 2023-2025"
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = RGBColor(100, 100, 100)

        print("✅ Slide 6: Répartition par marché")

    def add_architecture_modele_slide(self):
        """Slide 7: Architecture du modèle"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Titre
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        tf = title_box.text_frame
        tf.text = "ARCHITECTURE DU MODÈLE"
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_TITLE
        p.alignment = PP_ALIGN.CENTER

        # Placeholder pour graphique
        placeholder = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.8), Inches(1.2),
            Inches(8.4), Inches(5.5)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
        placeholder.line.color.rgb = RGBColor(200, 200, 200)

        # Texte dans placeholder
        text_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8.4), Inches(1))
        tf = text_box.text_frame
        tf.text = "📊 INSÉRER ICI: 05_architecture_modele.png"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(100, 100, 100)

        # Note
        note_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.4))
        tf = note_box.text_frame
        tf.text = "3 Piliers + Couche IA + Couche Décisionnelle (2 règles métier)"
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = RGBColor(100, 100, 100)

        print("✅ Slide 7: Architecture du modèle")

    def add_architecture_detail_slide(self):
        """Slide 8: Détail de l'architecture"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])  # Title and Content

        # Titre
        title = slide.shapes.title
        title.text = "DÉTAIL DE L'ARCHITECTURE"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.COLOR_TITLE

        # Contenu
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5.2))
        tf = content_box.text_frame
        tf.word_wrap = True

        # Piliers
        p = tf.add_paragraph()
        p.text = "🔵 PILIER 1: Type de Réclamation"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(52, 152, 219)
        p.space_after = Pt(6)

        p = tf.add_paragraph()
        p.text = "• Famille Produit  • Catégorie  • Sous-catégorie"
        p.font.size = Pt(16)
        p.level = 1
        p.space_after = Pt(12)

        p = tf.add_paragraph()
        p.text = "🔴 PILIER 2: Risque"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(231, 76, 60)
        p.space_after = Pt(6)

        p = tf.add_paragraph()
        p.text = "• Montant demandé  • Délai estimé  • Ratio Montant/PNB"
        p.font.size = Pt(16)
        p.level = 1
        p.space_after = Pt(12)

        p = tf.add_paragraph()
        p.text = "🟢 PILIER 3: Signalétique Client"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(46, 204, 113)
        p.space_after = Pt(6)

        p = tf.add_paragraph()
        p.text = "• PNB cumulé  • Ancienneté  • Segment  • Marché"
        p.font.size = Pt(16)
        p.level = 1
        p.space_after = Pt(18)

        # Couche analytique
        p = tf.add_paragraph()
        p.text = "🧠 COUCHE ANALYTIQUE"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(155, 89, 182)
        p.space_after = Pt(6)

        p = tf.add_paragraph()
        p.text = "• Modèles IA: XGBoost & CatBoost"
        p.font.size = Pt(16)
        p.level = 1

        p = tf.add_paragraph()
        p.text = "• Optimisation Optuna: Attribution automatique des POIDS optimaux"
        p.font.size = Pt(16)
        p.level = 1
        p.space_after = Pt(18)

        # Couche décisionnelle
        p = tf.add_paragraph()
        p.text = "⚙️ COUCHE DÉCISIONNELLE"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(243, 156, 18)
        p.space_after = Pt(6)

        p = tf.add_paragraph()
        p.text = "1️⃣ Décision Modèle: 3 zones (Rejet Auto | Audit Humain | Validation Auto)"
        p.font.size = Pt(16)
        p.level = 1

        p = tf.add_paragraph()
        p.text = "2️⃣ Règle métier #1: Maximum 1 validation par client par an"
        p.font.size = Pt(16)
        p.level = 1

        p = tf.add_paragraph()
        p.text = "3️⃣ Règle métier #2: Montant validé ≤ PNB année dernière"
        p.font.size = Pt(16)
        p.level = 1

        print("✅ Slide 8: Détail de l'architecture")

    def add_resultats_2025_slide(self):
        """Slide 9: Résultats 2025"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Titre
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        tf = title_box.text_frame
        tf.text = "RÉSULTATS 2025 & CALCUL DU GAIN"
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_TITLE
        p.alignment = PP_ALIGN.CENTER

        # Placeholder pour graphique
        placeholder = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.8), Inches(1.2),
            Inches(8.4), Inches(5.5)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
        placeholder.line.color.rgb = RGBColor(200, 200, 200)

        # Texte dans placeholder
        text_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8.4), Inches(1))
        tf = text_box.text_frame
        tf.text = "📊 INSÉRER ICI: 06_resultats_2025_gain.png"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(100, 100, 100)

        # Note
        note_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.4))
        tf = note_box.text_frame
        tf.text = "Points clés: • Taux d'automatisation  • Gain financier  • ETP libérés  • ROI"
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = RGBColor(100, 100, 100)

        print("✅ Slide 9: Résultats 2025")

    def add_benefices_slide(self):
        """Slide 10: Bénéfices"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])  # Title and Content

        # Titre
        title = slide.shapes.title
        title.text = "BÉNÉFICES DE L'OPÉRATIONNALISATION"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.COLOR_TITLE

        # Contenu en 2 colonnes
        # Colonne gauche - Bénéfices quantifiables
        left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4.2), Inches(5))
        tf = left_box.text_frame
        tf.word_wrap = True

        p = tf.add_paragraph()
        p.text = "💰 BÉNÉFICES QUANTIFIABLES"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_ACCENT
        p.space_after = Pt(10)

        benefits_quant = [
            ("Gain financier", "Réduction des coûts de traitement manuel"),
            ("Gain temps", "ETP libérés pour tâches à valeur ajoutée"),
            ("Productivité", "Augmentation significative du débit"),
            ("Réduction délais", "Traitement instantané vs manuel"),
            ("Cohérence", "Décisions standardisées et objectives")
        ]

        for title_text, desc in benefits_quant:
            p = tf.add_paragraph()
            p.text = f"✓ {title_text}"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.COLOR_SUBTITLE
            p.space_after = Pt(3)

            p = tf.add_paragraph()
            p.text = desc
            p.font.size = Pt(13)
            p.level = 1
            p.space_after = Pt(8)

        # Colonne droite - Bénéfices qualitatifs
        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.2), Inches(5))
        tf = right_box.text_frame
        tf.word_wrap = True

        p = tf.add_paragraph()
        p.text = "🎯 BÉNÉFICES QUALITATIFS"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_ACCENT
        p.space_after = Pt(10)

        benefits_qual = [
            ("Satisfaction client", "Réponses rapides et cohérentes"),
            ("Traçabilité", "Historique complet des décisions"),
            ("Transparence", "Explicabilité des choix du modèle"),
            ("Réduction erreurs", "Moins d'erreurs humaines"),
            ("Amélioration continue", "Monitoring et ajustements réguliers")
        ]

        for title_text, desc in benefits_qual:
            p = tf.add_paragraph()
            p.text = f"✓ {title_text}"
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.COLOR_SUBTITLE
            p.space_after = Pt(3)

            p = tf.add_paragraph()
            p.text = desc
            p.font.size = Pt(13)
            p.level = 1
            p.space_after = Pt(8)

        print("✅ Slide 10: Bénéfices")

    def add_recommandations_slide(self):
        """Slide 11: Recommandations"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])  # Title and Content

        # Titre
        title = slide.shapes.title
        title.text = "RECOMMANDATIONS & PROCHAINES ÉTAPES"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.COLOR_TITLE

        # Contenu
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.8))
        tf = content_box.text_frame
        tf.word_wrap = True

        recommandations = [
            ("1. DÉPLOIEMENT EN PRODUCTION", [
                "Mise en production du modèle validé",
                "Formation des équipes opérationnelles",
                "Documentation complète et guides utilisateurs"
            ]),
            ("2. MONITORING CONTINU", [
                "Tableau de bord de suivi quotidien",
                "Alertes automatiques sur dérives",
                "Rapports mensuels de performance"
            ]),
            ("3. GOUVERNANCE", [
                "Comité de pilotage trimestriel",
                "Revue des règles métier (semestrielle)",
                "Ré-entraînement annuel du modèle"
            ]),
            ("4. AMÉLIORATION CONTINUE", [
                "Collecte feedback utilisateurs",
                "Analyse des cas en audit humain",
                "Ajustements des seuils si nécessaire"
            ])
        ]

        for title_text, items in recommandations:
            p = tf.add_paragraph()
            p.text = title_text
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.COLOR_ACCENT
            p.space_after = Pt(6)

            for item in items:
                p = tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(14)
                p.level = 1
                p.space_after = Pt(3)

            p = tf.add_paragraph()
            p.text = ""
            p.space_after = Pt(8)

        print("✅ Slide 11: Recommandations")

    def add_conclusion_slide(self):
        """Slide 12: Conclusion"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Titre
        title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
        tf = title_box.text_frame
        tf.text = "CONCLUSION"
        p = tf.paragraphs[0]
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_TITLE
        p.alignment = PP_ALIGN.CENTER

        # Messages clés
        content_box = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(7), Inches(3))
        tf = content_box.text_frame
        tf.word_wrap = True

        messages = [
            "✓ Modèle robuste et performant validé sur 2024-2025",
            "✓ Architecture intelligente combinant IA et règles métier",
            "✓ Gains financiers et temps significatifs démontrés",
            "✓ Prêt pour déploiement en production immédiat",
            "✓ ROI positif dès la première année"
        ]

        for msg in messages:
            p = tf.add_paragraph()
            p.text = msg
            p.font.size = Pt(20)
            p.font.color.rgb = self.COLOR_SUBTITLE
            p.space_after = Pt(15)

        # Call to action
        cta_box = slide.shapes.add_textbox(Inches(2), Inches(6), Inches(6), Inches(0.8))
        tf = cta_box.text_frame
        tf.text = "🚀 GO/NO-GO pour le déploiement"
        p = tf.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.COLOR_ACCENT
        p.alignment = PP_ALIGN.CENTER

        print("✅ Slide 12: Conclusion")

    def add_questions_slide(self):
        """Slide 13: Questions"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank

        # Fond coloré
        background = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(10), Inches(7.5)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLOR_TITLE
        background.line.fill.background()

        # Texte Questions
        text_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
        tf = text_box.text_frame
        tf.text = "QUESTIONS ?"
        p = tf.paragraphs[0]
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        print("✅ Slide 13: Questions")

    def build(self):
        """Construire la présentation complète"""
        print("\n" + "="*80)
        print("🎨 GÉNÉRATION DE LA PRÉSENTATION POWERPOINT")
        print("="*80)

        # Slides
        self.add_title_slide()
        self.add_agenda_slide()

        # Section I: État des lieux
        self.add_section_slide("I", "ÉTAT DES LIEUX")
        self.add_evolution_volume_slide()
        self.add_fondee_vs_non_fondee_slide()
        self.add_repartition_famille_slide()
        self.add_repartition_marche_slide()

        # Section II: Modèle
        self.add_section_slide("II", "PRÉSENTATION DU MODÈLE")
        self.add_architecture_modele_slide()
        self.add_architecture_detail_slide()

        # Section III: Résultats
        self.add_section_slide("III", "RÉSULTATS & GAINS")
        self.add_resultats_2025_slide()
        self.add_benefices_slide()
        self.add_recommandations_slide()

        # Conclusion
        self.add_conclusion_slide()
        self.add_questions_slide()

        # Sauvegarder
        self.prs.save(self.output_file)

        print("\n" + "="*80)
        print("✅ PRÉSENTATION GÉNÉRÉE")
        print("="*80)
        print(f"\n📄 Fichier: {self.output_file}")
        print(f"📊 Nombre de slides: {len(self.prs.slides)}")
        print("\n💡 PROCHAINES ÉTAPES:")
        print("   1. Exécutez: python ml_pipeline_v2/generate_presentation_visuals.py")
        print("   2. Ouvrez la présentation PowerPoint")
        print("   3. Insérez les 6 graphiques PNG générés aux emplacements marqués")


def main():
    parser = argparse.ArgumentParser(description='Générer présentation PowerPoint')
    parser.add_argument('--output', type=str, default='outputs/presentation/presentation_scoring.pptx',
                       help='Nom du fichier de sortie')

    args = parser.parse_args()

    # Créer dossier de sortie
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Générer présentation
    builder = PresentationBuilder(output_path)
    builder.build()


if __name__ == '__main__':
    main()
