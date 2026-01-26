#!/usr/bin/env python3
"""
Générateur de Template PowerPoint - Attijariwafa Bank
Crée un template PowerPoint professionnel avec le logo de l'entreprise
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import argparse

# Couleurs Attijariwafa Bank
COLOR_ORANGE = RGBColor(241, 142, 51)  # Orange principal
COLOR_DARK_ORANGE = RGBColor(217, 93, 24)  # Orange foncé
COLOR_YELLOW = RGBColor(255, 195, 0)  # Jaune
COLOR_DARK_BLUE = RGBColor(0, 48, 87)  # Bleu foncé
COLOR_LIGHT_GRAY = RGBColor(242, 242, 242)  # Gris clair pour backgrounds
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_BLACK = RGBColor(0, 0, 0)
COLOR_TEXT_GRAY = RGBColor(64, 64, 64)

def add_logo(slide, logo_path, left=Inches(8.5), top=Inches(0.3), width=Inches(1.2)):
    """Ajoute le logo sur une slide si le fichier existe"""
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, left, top, width=width)
        return True
    return False

def add_footer(slide, text="", font_size=9):
    """Ajoute un pied de page à la slide"""
    left = Inches(0.5)
    top = Inches(7.2)
    width = Inches(9)
    height = Inches(0.3)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = text

    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = COLOR_TEXT_GRAY
    paragraph.alignment = PP_ALIGN.CENTER

def create_title_slide(prs, logo_path):
    """Crée une slide de titre"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Bannière orange en haut
    shapes = slide.shapes
    banner = shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(1.5)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = COLOR_ORANGE
    banner.line.fill.background()

    # Logo
    add_logo(slide, logo_path, left=Inches(8.3), top=Inches(0.2), width=Inches(1.4))

    # Titre principal
    title_box = shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Titre de la Présentation"
    title_frame.word_wrap = True

    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Sous-titre
    subtitle_box = shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Sous-titre ou description"

    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_TEXT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # Date et auteur
    info_box = shapes.add_textbox(Inches(1), Inches(6), Inches(8), Inches(0.5))
    info_frame = info_box.text_frame
    info_frame.text = "Date | Présenté par: [Votre Nom]"

    p = info_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_TEXT_GRAY
    p.alignment = PP_ALIGN.CENTER

    return slide

def create_section_divider_slide(prs, logo_path):
    """Crée une slide de séparation de section"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background orange
    shapes = slide.shapes
    bg = shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_ORANGE
    bg.line.fill.background()

    # Logo en blanc si possible, sinon normal
    add_logo(slide, logo_path, left=Inches(8.3), top=Inches(0.2), width=Inches(1.4))

    # Numéro de section
    section_num_box = shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    section_frame = section_num_box.text_frame
    section_frame.text = "PARTIE 1"

    p = section_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Titre de section
    title_box = shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Titre de la Section"
    title_frame.word_wrap = True

    p = title_frame.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide

def create_content_slide(prs, logo_path):
    """Crée une slide de contenu standard"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    shapes = slide.shapes

    # Bannière de titre
    title_banner = shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    title_banner.fill.solid()
    title_banner.fill.fore_color.rgb = COLOR_ORANGE
    title_banner.line.fill.background()

    # Logo
    add_logo(slide, logo_path, left=Inches(8.5), top=Inches(0.15), width=Inches(1.2))

    # Titre
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(7.5), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Titre de la Slide"
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.LEFT

    # Zone de contenu
    content_box = shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    # Ajouter des bullet points d'exemple
    content_frame.text = "Point principal 1"
    p = content_frame.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_TEXT_GRAY
    p.level = 0

    for text in ["Point principal 2", "Point principal 3", "Point principal 4"]:
        p = content_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXT_GRAY
        p.level = 0

    add_footer(slide, "Attijariwafa Bank | Confidentiel")

    return slide

def create_two_column_slide(prs, logo_path):
    """Crée une slide avec 2 colonnes"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    shapes = slide.shapes

    # Bannière de titre
    title_banner = shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    title_banner.fill.solid()
    title_banner.fill.fore_color.rgb = COLOR_ORANGE
    title_banner.line.fill.background()

    # Logo
    add_logo(slide, logo_path, left=Inches(8.5), top=Inches(0.15), width=Inches(1.2))

    # Titre
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(7.5), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Titre - Comparaison 2 Colonnes"
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    # Colonne gauche
    left_box = shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    left_frame.text = "Colonne Gauche"

    p = left_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK_BLUE

    p = left_frame.add_paragraph()
    p.text = "• Point 1\n• Point 2\n• Point 3"
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_TEXT_GRAY

    # Colonne droite
    right_box = shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.5), Inches(5.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    right_frame.text = "Colonne Droite"

    p = right_frame.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK_BLUE

    p = right_frame.add_paragraph()
    p.text = "• Point 1\n• Point 2\n• Point 3"
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_TEXT_GRAY

    add_footer(slide, "Attijariwafa Bank | Confidentiel")

    return slide

def create_graph_slide(prs, logo_path):
    """Crée une slide pour insérer un graphique"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    shapes = slide.shapes

    # Bannière de titre
    title_banner = shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    title_banner.fill.solid()
    title_banner.fill.fore_color.rgb = COLOR_ORANGE
    title_banner.line.fill.background()

    # Logo
    add_logo(slide, logo_path, left=Inches(8.5), top=Inches(0.15), width=Inches(1.2))

    # Titre
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(7.5), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Résultats - Graphique"
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    # Zone réservée pour graphique
    graph_placeholder = shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(1.2),
        Inches(9), Inches(5.3)
    )
    graph_placeholder.fill.solid()
    graph_placeholder.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    graph_placeholder.line.color.rgb = COLOR_ORANGE
    graph_placeholder.line.width = Pt(2)

    # Texte dans la zone
    text_box = shapes.add_textbox(Inches(3.5), Inches(3.5), Inches(3), Inches(0.8))
    text_frame = text_box.text_frame
    text_frame.text = "Insérer le graphique ici\n(Supprimez ce texte)"

    p = text_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_TEXT_GRAY
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide, "Attijariwafa Bank | Confidentiel")

    return slide

def create_conclusion_slide(prs, logo_path):
    """Crée une slide de conclusion"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    shapes = slide.shapes

    # Background dégradé (simulé avec orange clair)
    bg = shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    bg.line.fill.background()

    # Logo
    add_logo(slide, logo_path, left=Inches(8.3), top=Inches(0.2), width=Inches(1.4))

    # Bannière orange pour le titre
    banner = shapes.add_shape(
        1,  # Rectangle
        Inches(1.5), Inches(2.5),
        Inches(7), Inches(1.2)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = COLOR_ORANGE
    banner.line.fill.background()

    # Titre
    title_box = shapes.add_textbox(Inches(1.5), Inches(2.6), Inches(7), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Merci de votre attention"
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Contact
    contact_box = shapes.add_textbox(Inches(2), Inches(4.5), Inches(6), Inches(1.5))
    contact_frame = contact_box.text_frame
    contact_frame.word_wrap = True

    contact_frame.text = "Questions?"
    p = contact_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    p = contact_frame.add_paragraph()
    p.text = "\nContact: votre.email@attijariwafabank.com"
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_TEXT_GRAY
    p.alignment = PP_ALIGN.CENTER

    return slide

def main():
    parser = argparse.ArgumentParser(description='Génère un template PowerPoint Attijariwafa Bank')
    parser.add_argument('--logo', type=str,
                       default='assets/attijariwafa_logo.png',
                       help='Chemin vers le logo Attijariwafa Bank')
    parser.add_argument('--output', type=str,
                       default='outputs/Template_Attijariwafa.pptx',
                       help='Chemin de sortie pour le template')

    args = parser.parse_args()

    # Créer le dossier de sortie si nécessaire
    output_dir = os.path.dirname(args.output)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir)

    # Vérifier si le logo existe
    logo_exists = os.path.exists(args.logo)
    if not logo_exists:
        print(f"⚠️  Logo non trouvé: {args.logo}")
        print(f"💡 Le template sera créé sans logo. Ajoutez le logo manuellement ou")
        print(f"   placez le fichier logo à: {args.logo}")
        print()

    # Créer la présentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("📊 Génération du template PowerPoint Attijariwafa Bank...")
    print()

    # 1. Slide de titre
    print("✓ Slide 1: Page de titre")
    create_title_slide(prs, args.logo)

    # 2. Slide séparateur de section
    print("✓ Slide 2: Séparateur de section")
    create_section_divider_slide(prs, args.logo)

    # 3. Slide de contenu standard
    print("✓ Slide 3: Contenu standard (bullet points)")
    create_content_slide(prs, args.logo)

    # 4. Slide 2 colonnes
    print("✓ Slide 4: Deux colonnes (comparaison)")
    create_two_column_slide(prs, args.logo)

    # 5. Slide pour graphique
    print("✓ Slide 5: Zone graphique")
    create_graph_slide(prs, args.logo)

    # 6. Slide de conclusion
    print("✓ Slide 6: Conclusion / Questions")
    create_conclusion_slide(prs, args.logo)

    # Sauvegarder
    prs.save(args.output)

    print()
    print(f"✅ Template créé avec succès: {args.output}")
    print()
    print("📋 Le template contient 6 slides types:")
    print("   1. Page de titre")
    print("   2. Séparateur de section (fond orange)")
    print("   3. Contenu standard avec bullet points")
    print("   4. Deux colonnes (comparaison)")
    print("   5. Zone pour insérer graphiques")
    print("   6. Conclusion / Questions")
    print()
    print("🎨 Couleurs Attijariwafa Bank appliquées:")
    print("   - Orange principal: RGB(241, 142, 51)")
    print("   - Orange foncé: RGB(217, 93, 24)")
    print("   - Jaune: RGB(255, 195, 0)")
    print("   - Bleu foncé: RGB(0, 48, 87)")
    print()

    if logo_exists:
        print("✅ Logo intégré sur toutes les slides")
    else:
        print("⚠️  Pour ajouter le logo:")
        print(f"   1. Placez votre logo à: {args.logo}")
        print(f"   2. Relancez le script")
        print("   OU")
        print("   3. Ajoutez-le manuellement dans PowerPoint")
    print()
    print("💡 Dupliquez les slides pour créer votre présentation!")

if __name__ == "__main__":
    main()
